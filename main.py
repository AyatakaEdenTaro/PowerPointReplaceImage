import shutil
from pptx import Presentation

# 画像ファイルの読み込み
report_img_path = "./after/report.png"

# PowerPointファイル読み込み
prs = Presentation("./レポート.pptx")

# 画像の追加と配置元シェイプの削除
for i, sld in enumerate(prs.slides, start=1):
    for shape in sld.shapes:
        if(shape.text=="ここに画像"):
            sld.shapes.add_picture(report_img_path,
            shape.left,
            shape.top,
            width=shape.width,
            height=shape.height)
            sp = shape.element
            sp.getparent().remove(sp)

prs.save("./コピーレポート.pptx")